import os
import time

from pptx import Presentation
from pptx.util import Inches
import pyautogui
import tkinter as tk
from tkinter import filedialog


class PyPPTXShot(object):
    """PyPPTXShot main class."""

    def __init__(self, root) -> None:
        self.root = root
        self.root.title("PyPPTXShot - A Python Screenshotting PPTX Generator")
        self.root.resizable(0, 0)

        self.rect = None
        self.x = 0
        self.y = 0
        self.start_x = 0
        self.start_y = 0
        self.current_x = 0
        self.current_y = 0
        self.do_fail_safe = False
        self.MAXIMUM_SIZE = 51206400
        self.f2_text = "Use Current Coords. (F2 for instant)"

        # Defaults for screenshotter.
        self.click_delay_seconds = 0.1
        self.timer_delay_seconds = 3
        self.slide_multiplier = 1
        self.screenshots_per_slide = 4
        self.columns = 3
        self.scaling_factor = 100

        # Defaults for generated PPTX file.
        self.num_of_slides = 20
        self.filepath = os.path.expanduser("~\Desktop")
        self.filename = "pypptxshot_output"

        # Filename.
        self.use_filename = tk.Label(self.root, text="Filename:")
        self.use_filename.grid(row=1, column=0)
        self.filename_entry = tk.Entry(self.root)
        self.filename_entry.grid(row=1, column=1)
        self.filename_entry.insert(0, self.filename)

        # Screenshot Multiplier.
        self.multiply_screenshots = tk.Label(
            self.root, text="Screenshot Multiplier:"
        )
        self.multiply_screenshots.grid(row=2, column=0)
        self.multiply_screenshots_entry = tk.Entry(self.root)
        self.multiply_screenshots_entry.grid(row=2, column=1)
        self.multiply_screenshots_entry.insert(0, self.slide_multiplier)

        # Number of Screenshots.
        self.num_screenshots = tk.Label(self.root, text="# of Screenshot(s)")
        self.num_screenshots.grid(row=3, column=0)
        self.num_screenshots_entry = tk.Entry(self.root)
        self.num_screenshots_entry.grid(row=3, column=1)
        self.num_screenshots_entry.insert(0, self.num_of_slides)

        # Screenshots per Slide.
        self.ss_per_slide = tk.Label(
            self.root, text="Screenshot(s) per Slide:"
        )
        self.ss_per_slide.grid(row=4, column=0)
        self.ss_per_slide_entry = tk.Entry(self.root)
        self.ss_per_slide_entry.grid(row=4, column=1)
        self.ss_per_slide_entry.insert(0, self.screenshots_per_slide)

        # Maximum columns per slide.
        self.column_per_slide = tk.Label(
            self.root, text="Column(s) per Slide:"
        )
        self.column_per_slide.grid(row=5, column=0)
        self.column_per_slide_entry = tk.Entry(self.root)
        self.column_per_slide_entry.grid(row=5, column=1)
        self.column_per_slide_entry.insert(0, self.columns)

        # Scaling of each screenshot to be placed in the file.
        self.scaling = tk.Label(self.root, text="Screenshot Scaling (%):")
        self.scaling.grid(row=6, column=0)
        self.scaling_entry = tk.Entry(self.root)
        self.scaling_entry.grid(row=6, column=1)
        self.scaling_entry.insert(0, self.scaling_factor)

        # Save location of PPTX file.
        self.save_directory = tk.Label(self.root, text="Save to Directory:")
        self.save_directory.grid(row=7, column=0)
        self.save_directory_entry = tk.Entry(self.root)
        self.save_directory_entry.grid(row=7, column=1, columnspan=3, ipadx=80)
        self.save_directory_entry.insert(0, r"{}".format(self.filepath))
        self.save_directory_button = tk.Button(
            self.root, text="Browse...", command=self.ask_directory
        )
        self.save_directory_button.grid(row=7, column=4, ipadx=50)

        # Checkbox of manual XY coordinates.
        self.manual_xy = tk.Label(
            self.root,
            text="Auto-click XY",
        )
        self.manual_xy.grid(row=2, column=3)
        self.manual_xy_entry = tk.Entry(self.root)
        self.manual_xy_entry.grid(row=2, column=4)

        # PyAutoGUI based.
        self.current_xy = tk.Label(self.root, text="Mouse Current XY")
        self.current_xy.grid(row=1, column=3)
        self.current_xy_entry = tk.Entry(self.root)
        self.current_xy_entry.grid(row=1, column=4)

        # Click delay.
        self.click_delay = tk.Label(self.root, text="Click Delay (sec.)")
        self.click_delay.grid(row=3, column=3)
        self.click_delay_entry = tk.Entry(self.root)
        self.click_delay_entry.grid(row=3, column=4)
        self.click_delay_entry.insert(0, self.click_delay_seconds)

        # Checkbox for auto-click, default is false.
        self.status_do_autoclick = tk.IntVar(value=0)
        self.do_autoclick = tk.Checkbutton(
            self.root,
            text="Do auto-click",
            variable=self.status_do_autoclick,
            offvalue=0,
            onvalue=1,
        )
        self.do_autoclick.grid(row=4, column=3, columnspan=2, rowspan=2)

        # Checkbox if user wants to scale the slide size according to
        # the screenshot scale. Defaults to true.
        self.status_slide_scaled = tk.IntVar(value=1)
        self.slide_scaled = tk.Checkbutton(
            self.root,
            text="Adjust slide size to scaled screenshot",
            variable=self.status_slide_scaled,
            offvalue=0,
            onvalue=1,
        )
        self.slide_scaled.grid(row=5, column=3, columnspan=2, rowspan=2)

        # Manually get coordinates after a timer.
        self.get_coords_text = tk.StringVar()
        self.get_coords_text.set(self.f2_text)
        self.get_coords_button = tk.Button(
            self.root,
            width=3,
            command=self.change_coords,
            textvariable=self.get_coords_text,
        )
        self.get_coords_button.grid(row=8, column=3, columnspan=2, ipadx=180)

        # Button of screen capture.
        self.snip_button = tk.Button(
            self.root,
            width=3,
            command=self.create_screen_canvas,
            text="Begin Screen Capture...",
        )
        self.snip_button.grid(row=8, column=0, columnspan=2, ipadx=180)

        # Bind.
        self.root.bind("<F2>", lambda a=True: self.change_coords(a))

        self.root_screen = tk.Toplevel(self.root)
        self.root_screen.withdraw()
        self.picture_frame = tk.Frame(self.root_screen)
        self.picture_frame.pack(fill=tk.BOTH, expand=tk.YES)
        self.get_mouse_pos()

    def ask_directory(self) -> None:
        """
        Ask for user's desired directory on where to save the generated
        PPTX file.
        """

        directory_chosen = filedialog.askdirectory(
            title="Select folder for output"
        )
        self.save_directory_entry.delete(0, tk.END)
        self.save_directory_entry.insert(0, directory_chosen)

    def change_coords(self, instant=False) -> None:
        """
        Changes the coordinates to be auto-clicked.

        Waits for a few seconds (see self.timer_delay_seconds) before
        getting the coordinates of your pointer.
        """
        if not instant:
            for second in range(self.timer_delay_seconds):
                self.get_coords_text.set(
                    f"Align mouse in {3 - second} seconds..."
                )
                self.root.update()
                time.sleep(1)
        data = self.current_xy_entry.get()
        self.x_click, self.y_click = data.split(",")
        self.get_coords_text.set(self.f2_text)
        self.manual_xy_entry.delete(0, tk.END)
        self.manual_xy_entry.insert(0, data)

    def get_mouse_pos(self) -> None:
        """Updates mouse position until program exits."""

        self.current_xy_entry.delete(0, tk.END)
        self.current_xy_entry.insert(
            0, "{},{}".format(*self.root.winfo_pointerxy())
        )
        self.root.after(100, self.get_mouse_pos)

    def take_bounded_screenshot(self, x1, y1, x2, y2) -> None:
        """Takes a screenshot based on the coordinates given.

        Main function for placing and saving screenshots to PPTX file."""

        # Get mouse coords.
        manual_coords = self.manual_xy_entry.get()
        default_coords = (
            float(self.start_x + self.current_x) / 2,
            float(self.start_y + self.current_y) / 2,
        )
        if manual_coords == 1 or manual_coords != "":
            try:
                self.x_click, self.y_click = [
                    int(x) for x in manual_coords.split(",")
                ]
            except ValueError:
                print("Error in manual coords. Defaulted to center.")
                self.x_click, self.y_click = default_coords
        else:
            self.change_coords(instant=True)

        # Get number of slides.
        if self.num_screenshots_entry.get() != "":
            self.num_of_slides = int(self.num_screenshots_entry.get())

        # Multiply slides.
        self.num_of_slides *= int(self.multiply_screenshots_entry.get())

        # Initialize.
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slides = []

        # Update.
        error_message = (
            f"Not an integer. Defaulting to 4 for screenshots per slide.",
        )
        try:
            self.screenshots_per_slide = int(self.ss_per_slide_entry.get())
        except ValueError:
            print(error_message)
            self.screenshots_per_slide = 4

        # Update columns.
        try:
            self.columns = int(self.column_per_slide_entry.get())
        except ValueError:
            print("Invalid column limit. Defaulting to 2.")
            self.columns = 2
        print("Current number of columns per slide:", self.columns)

        # Update rows, if column number is odd increment 1 due to space issue.
        rows = int(self.screenshots_per_slide / self.columns)
        self.is_odd = False
        if self.screenshots_per_slide / self.columns != 1:
            self.is_odd = True
        if self.screenshots_per_slide % self.columns != 0 or rows < 1:
            rows += 1
        print("Current number of rows per slide:", rows)

        # Width and height is adjusted based on screenshot dimensions.
        # Inch value taken from:
        # https://www.unitconverters.net/typography/pixel-x-to-inch.htm
        INCH = 0.0104166667

        try:
            self.scaling_factor = float(self.scaling_entry.get())
        except ValueError:
            print("Error in scaling factor. Did you input a valid number?")
            print("Defaulting to 100%.")
            self.scaling_factor = 100

        if self.scaling_factor > 100:
            print(
                "Warning. Scaling is greater than 100%. Results may not be"
                " pleasant."
            )

        scale = self.scaling_factor / 100

        slide_width = (x2 * INCH) * self.columns
        slide_height = (y2 * INCH) * rows

        if self.status_slide_scaled.get():
            slide_width = slide_width * scale
            slide_height = slide_height * scale

        screenshot_width = (x2 * INCH) * scale
        screenshot_height = (y2 * INCH) * scale

        if Inches(slide_width) > self.MAXIMUM_SIZE:
            print(
                f"Slide width {slide_width} reached beyond maximum of"
                f" python-pptx. Reverting to {self.MAXIMUM_SIZE}."
            )
            prs.slide_width = self.MAXIMUM_SIZE
        else:
            prs.slide_width = Inches(slide_width)

        if Inches(slide_height) > self.MAXIMUM_SIZE:
            print(
                f"Slide height {slide_height} reached beyond maximum of"
                f" python-pptx. Reverting to {self.MAXIMUM_SIZE}."
            )
            prs.slide_height = self.MAXIMUM_SIZE
        else:
            prs.slide_height = Inches(slide_height)

        print(f"Slide size is {slide_width}x{slide_height} inches.")
        print(
            "Size for each screenshot in slide is",
            f"{screenshot_width}x{screenshot_height} inches.",
        )
        print("First cursor hold coordinates:", x1, y1)
        print("Last cursor hold coordinates:", x2, y2)
        print(
            "Calculated cursor midpoint to auto-click (if enabled):",
            self.x_click,
            self.y_click,
        )

        # Add the base slides.
        for _ in range(1, self.num_of_slides + 1, self.screenshots_per_slide):
            slides.append(prs.slides.add_slide(blank_slide_layout))
        print(
            f"{len(slides)} slides are made out of {self.num_of_slides}",
            f"configured slots via skipping by {self.screenshots_per_slide}.",
        )

        # Begin screenshotting. Save to a PNG (.png) file first so it
        # be aligned in the PPTX (.pptx) file.
        try:
            # Reverse.
            if self.is_odd:
                self.columns, rows = rows, self.columns

            left_inches = slide_width / rows
            top_inches = slide_height / self.columns

            for slide_number in range(len(slides)):
                current_shot = 1
                for column in range(self.columns):
                    for row in range(rows):
                        # Retrieve delay time and wait before clicking.
                        if self.status_do_autoclick.get() == 1:
                            time.sleep(float(self.click_delay_entry.get()))
                            pyautogui.click(x=self.x_click, y=self.y_click)

                        if current_shot <= self.screenshots_per_slide:
                            # Take a shot, save, then deal with the placement.
                            print(
                                f"[{column}, {row}] CURRENTLY SCREENSHOTTING"
                                f" FOR SLIDE #{slide_number + 1}."
                            )

                            img_path = f"ppg_cache.png"
                            pyautogui.screenshot(
                                img_path, region=(x1, y1, x2, y2)
                            )

                            top_final = top_inches * column
                            left_final = left_inches * row
                            print("TOP, LEFT:", (top_final, left_final))

                            slides[slide_number].shapes.add_picture(
                                img_path,
                                top=Inches(top_final),
                                left=Inches(left_final),
                                width=Inches(screenshot_width),
                                height=Inches(screenshot_height),
                            )

                            # Remove cache.
                            os.remove(img_path)

                            current_shot += 1

        except pyautogui.FailSafeException:
            self.do_fail_safe = True

        # Do a file iteration check for no accidental overwriting.
        counter = 0
        while (
            os.path.exists(
                os.path.join(self.filepath, self.filename + f"-{counter}.pptx")
            )
            and not self.do_fail_safe
        ):
            print(f"File No. #{counter} detected. Changing file enumeration.")
            counter += 1
        final_savepath = os.path.join(
            self.filepath, self.filename + f"-{counter}.pptx"
        )

        # Save and log.
        if not self.do_fail_safe:
            prs.save(final_savepath)
            print(
                f"Screenshotting complete. File is saved in: {final_savepath}"
            )
        else:
            self.do_fail_safe = False
            print("Fail safe detected. Cancelled operation.")

    def create_screen_canvas(self) -> None:
        """
        Generates an area using tk.Canvas that can be used to select a
        rectangular space (click and hold) on where will the app would
        screenshot. Will begin screenshotting as soon as the mouse is
        released. Press the right mouse button (right-click) to cancel.
        """
        self.root_screen.deiconify()
        self.root.withdraw()

        self.screenCanvas = tk.Canvas(
            self.picture_frame, cursor="cross", bg="grey11"
        )
        self.screenCanvas.pack(fill=tk.BOTH, expand=tk.YES)

        self.screenCanvas.bind("<ButtonPress-3>", self.quick_cancel)
        self.screenCanvas.bind("<ButtonPress-1>", self.on_button_press)
        self.screenCanvas.bind("<B1-Motion>", self.on_button_hold)
        self.screenCanvas.bind("<ButtonRelease-1>", self.on_button_release)

        self.root_screen.attributes("-fullscreen", True)
        self.root_screen.attributes("-alpha", 0.3)
        self.root_screen.lift()
        self.root_screen.attributes("-topmost", True)

    def quick_cancel(self, event):
        """
        Handles event if user decided to exit screenshot mode.
        """
        self.exit_screenshot_mode()
        self.root.deiconify()
        return event

    def on_button_release(self, event):
        """
        Handles event if left mouse button is released during
        create_screen_canvas.
        """

        self.exit_screenshot_mode()
        if self.start_x <= self.current_x and self.start_y <= self.current_y:
            self.take_bounded_screenshot(
                self.start_x,
                self.start_y,
                self.current_x - self.start_x,
                self.current_y - self.start_y,
            )

        elif self.start_x >= self.current_x and self.start_y <= self.current_y:
            self.take_bounded_screenshot(
                self.current_x,
                self.start_y,
                self.start_x - self.current_x,
                self.current_y - self.start_y,
            )

        elif self.start_x <= self.current_x and self.start_y >= self.current_y:
            self.take_bounded_screenshot(
                self.start_x,
                self.current_y,
                self.current_x - self.start_x,
                self.start_y - self.current_y,
            )

        elif self.start_x >= self.current_x and self.start_y >= self.current_y:
            self.take_bounded_screenshot(
                self.current_x,
                self.current_y,
                self.start_x - self.current_x,
                self.start_y - self.current_y,
            )
        self.root.deiconify()
        return event

    def on_button_press(self, event) -> None:
        """
        Handles event of the first click of mouse during
        create_screen_canvas.
        """

        # Save mouse at the start.
        self.start_x = self.screenCanvas.canvasx(event.x)
        self.start_y = self.screenCanvas.canvasy(event.y)
        self.rect = self.screenCanvas.create_rectangle(
            self.x, self.y, 1, 1, width=3, fill="blue"
        )

    def on_button_hold(self, event) -> None:
        """
        Handles event of click and hold feature of mouse during
        create_screen_canvas.
        """

        self.current_x, self.current_y = (event.x, event.y)
        # Expand rectangle when left mouse button is hold.
        self.screenCanvas.coords(
            self.rect,
            self.start_x,
            self.start_y,
            self.current_x,
            self.current_y,
        )

    def exit_screenshot_mode(self) -> None:
        """Function for exiting create_screen_canvas."""

        print("Screenshot mode exited.")
        self.screenCanvas.destroy()
        self.root_screen.withdraw()

    def exit_application(self) -> None:
        """Quits PyPPTXShot."""

        print("PyPPTXShot is quitting.")
        self.root.quit()

if __name__ == "__main__":
    app = tk.Tk()
    pptx_shot = PyPPTXShot(app)
    app.mainloop()
